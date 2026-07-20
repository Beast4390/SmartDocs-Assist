import os
import re
import uuid
import unicodedata
from datetime import datetime

# Import document parsing libraries
try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    import docx  # python-docx
except ImportError:
    docx = None

try:
    import pptx  # python-pptx
except ImportError:
    pptx = None


class PasswordProtectedPDFException(Exception):
    """Exception raised when a PDF file is encrypted/password-protected."""
    pass


class ExtractionFailureException(Exception):
    """Exception raised when document text extraction fails."""
    pass


class DocumentProcessor:
    """
    Enterprise-grade Document Processing Service.
    Handles secure validation, extraction, cleaning, and metadata compilation 
    for PDF, DOCX, and PPTX formats, entirely offline and air-gapped.
    """
    
    def __init__(self, upload_folder: str = "uploads", processed_folder: str = "processed", metadata_folder: str = "metadata"):
        self.upload_folder = upload_folder
        self.processed_folder = processed_folder
        self.metadata_folder = metadata_folder
        
        # Ensure workspace directories exist
        os.makedirs(self.upload_folder, exist_ok=True)
        os.makedirs(self.processed_folder, exist_ok=True)
        os.makedirs(self.metadata_folder, exist_ok=True)

        self.allowed_extensions = {'pdf', 'docx', 'pptx'}
        self.max_file_size = 50 * 1024 * 1024  # 50MB as specified in user guidelines and UI

    def validate_file(self, filename: str, file_size: int) -> tuple[bool, str]:
        """
        Validate that the uploaded file complies with type and size limits.
        """
        if not filename or '.' not in filename:
            return False, "File must have a valid name and extension."
            
        ext = filename.rsplit('.', 1)[1].lower()
        if ext not in self.allowed_extensions:
            return False, f"Unsupported file extension '.{ext}'. Allowed formats: {', '.join(sorted(list(self.allowed_extensions))).upper()}."
            
        if file_size > self.max_file_size:
            return False, f"File size exceeds maximum allowed limit of 50MB. (Uploaded: {file_size / (1024 * 1024):.2f}MB)"
            
        return True, ""

    def save_file(self, file_storage, filename: str) -> str:
        """
        Saves the raw uploaded file storage stream to the uploads directory.
        """
        from werkzeug.utils import secure_filename
        
        safe_name = secure_filename(filename)
        # Create unique filename base to prevent collisions on identical uploads
        name_parts = os.path.splitext(safe_name)
        ext = name_parts[1].lower()
        clean_base = re.sub(r'[^a-zA-Z0-9_\-]', '_', name_parts[0])[:50]
        unique_id = uuid.uuid4().hex[:8]
        unique_name = f"{clean_base}_{unique_id}{ext}"
        
        save_path = os.path.join(self.upload_folder, unique_name)
        file_storage.save(save_path)
        return unique_name

    def clean_text(self, text: str) -> str:
        """
        Applies standard corporate normalization rules on raw text blocks:
        - Normalizes Unicode representations to NFKC
        - Collapses duplicate/excess blank lines (retaining structure)
        - Cleans tabulator/carriage-return variations
        - Preserves punctuation and paragraph breaks
        """
        if not text:
            return ""
            
        # Normalize Unicode characters
        text = unicodedata.normalize('NFKC', text)
        
        # Standardize carriage returns to newlines
        text = text.replace('\r\n', '\n').replace('\r', '\n')
        
        # Replace vertical spaces, tabs or weird spaces with single standard spaces on each line
        lines = []
        for line in text.split('\n'):
            # Strip multiple trailing/leading spaces, replace multiple inline spaces/tabs with single space
            cleaned_line = re.sub(r'[ \t]+', ' ', line).strip()
            lines.append(cleaned_line)
            
        # Collapse multiple consecutive blank lines into at most a single blank line (double newline)
        collapsed_text = '\n'.join(lines)
        collapsed_text = re.sub(r'\n{3,}', '\n\n', collapsed_text)
        
        return collapsed_text.strip()

    def extract_pdf(self, file_path: str) -> tuple[str, int]:
        """
        Extracts content from PDF files using PyMuPDF (fitz) safely,
        skipping corrupted pages and failing gracefully on encrypted PDFs.
        """
        if fitz is None:
            raise ImportError("PyMuPDF (fitz) is not installed in the current environment.")
            
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"PDF file not found at path: {file_path}")
            
        try:
            doc = fitz.open(file_path)
        except Exception as e:
            raise ExtractionFailureException(f"Failed to open PDF document. It may be corrupted. Error: {str(e)}")
            
        if doc.is_encrypted:
            raise PasswordProtectedPDFException("Password-protected or encrypted PDF files are not supported.")
            
        pages_text = []
        total_pages = doc.page_count
        
        for i in range(total_pages):
            try:
                page = doc.load_page(i)
                text = page.get_text("text")
                if text:
                    pages_text.append(text)
            except Exception as page_err:
                # Log page skip and continue to extract remaining text
                print(f"[Warning] Skipped corrupted page {i+1} in {file_path}: {str(page_err)}")
                continue
                
        doc.close()
        combined_text = "\n\n".join(pages_text)
        return combined_text, total_pages

    def extract_docx(self, file_path: str) -> tuple[str, int]:
        """
        Extracts paragraphs and tables from DOCX using python-docx safely.
        Approximates a page count based on average word count.
        """
        if docx is None:
            raise ImportError("python-docx is not installed in the current environment.")
            
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"DOCX file not found at path: {file_path}")
            
        try:
            doc = docx.Document(file_path)
        except Exception as e:
            raise ExtractionFailureException(f"Failed to open DOCX document. It may be corrupted. Error: {str(e)}")
            
        elements_text = []
        
        # 1. Extract headings and normal paragraphs
        for p in doc.paragraphs:
            text_val = p.text.strip()
            if text_val:
                elements_text.append(text_val)
                
        # 2. Extract tabular contents
        for table in doc.tables:
            for row in table.rows:
                row_cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_cells:
                    elements_text.append(" | ".join(row_cells))
                    
        combined_text = "\n\n".join(elements_text)
        
        # Approximate standard physical page count for DOCX (approx 500 words/page, minimum 1)
        word_count = len(combined_text.split())
        approx_pages = max(1, (word_count // 500) + 1)
        
        return combined_text, approx_pages

    def extract_pptx(self, file_path: str) -> tuple[str, int]:
        """
        Extracts titles, text boxes, bullets, and slide notes from PPTX using python-pptx safely.
        """
        if pptx is None:
            raise ImportError("python-pptx is not installed in the current environment.")
            
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"PPTX file not found at path: {file_path}")
            
        try:
            prs = pptx.Presentation(file_path)
        except Exception as e:
            raise ExtractionFailureException(f"Failed to open PPTX presentation. It may be corrupted. Error: {str(e)}")
            
        slides_text = []
        slide_count = len(prs.slides)
        
        for index, slide in enumerate(prs.slides):
            slide_number = index + 1
            slide_elements = [f"--- [Slide {slide_number}] ---"]
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_elements.append(shape.text.strip())
                    
            # Extract notes from notes slides if present
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_elements.append(f"[Speaker Notes]: {notes_text}")
                    
            slides_text.append("\n".join(slide_elements))
            
        combined_text = "\n\n".join(slides_text)
        return combined_text, slide_count

    def create_metadata(self, filename: str, original_filename: str, ext: str, file_size: int, text: str, page_or_slide_count: int, status: str = "processed") -> dict:
        """
        Compiles structural and statistics metadata for the processed document.
        """
        word_count = len(text.split()) if text else 0
        char_count = len(text) if text else 0
        
        metadata = {
            "document_id": str(uuid.uuid4()),
            "filename": filename,
            "original_filename": original_filename,
            "extension": ext.upper(),
            "upload_date": datetime.utcnow().strftime("%Y-%m-%dT%H:%M:%SZ"),
            "file_size": file_size,
            "pages": page_or_slide_count,
            "words": word_count,
            "characters": char_count,
            "status": status
        }
        return metadata

    def process_document(self, file_path: str, original_filename: str) -> dict:
        """
        Main orchestrator pipeline for document ingestion.
        Loads file, identifies extension, extracts text, normalizes/cleans text, 
        compiles JSON metadata, and writes output files safely to disk.
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"Target file for processing not found at path: {file_path}")
            
        filename = os.path.basename(file_path)
        ext = filename.rsplit('.', 1)[1].lower()
        file_size = os.path.getsize(file_path)
        
        # 1. Extraction based on type
        if ext == 'pdf':
            raw_text, page_count = self.extract_pdf(file_path)
        elif ext == 'docx':
            raw_text, page_count = self.extract_docx(file_path)
        elif ext == 'pptx':
            raw_text, page_count = self.extract_pptx(file_path)
        else:
            raise ValueError(f"Unmapped file extension: {ext}")
            
        # 2. Check for empty document extraction
        cleaned_text = self.clean_text(raw_text)
        if not cleaned_text:
            raise ValueError("Document yielded no extractable text content.")
            
        # 3. Create Metadata record
        metadata = self.create_metadata(
            filename=filename,
            original_filename=original_filename,
            ext=ext,
            file_size=file_size,
            text=cleaned_text,
            page_or_slide_count=page_count,
            status="processed"
        )
        
        # 4. Save Processed Text to processed/ folder
        txt_filename = f"{os.path.splitext(filename)[0]}.txt"
        txt_path = os.path.join(self.processed_folder, txt_filename)
        with open(txt_path, "w", encoding="utf-8") as txt_file:
            txt_file.write(cleaned_text)
            
        # 5. Save Metadata as JSON to metadata/ folder
        json_filename = f"{os.path.splitext(filename)[0]}.json"
        json_path = os.path.join(self.metadata_folder, json_filename)
        import json
        with open(json_path, "w", encoding="utf-8") as json_file:
            json.dump(metadata, json_file, indent=4, ensure_ascii=False)
            
        return metadata
